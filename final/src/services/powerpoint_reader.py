import os
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging
from PIL import Image
import io

from ..types.models import SlideContent, PresentationData, SlideType

class PowerPointReader:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def read_presentation(self, file_path: str) -> PresentationData:
        """PowerPointファイルを読み込み、構造化データに変換"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = self._extract_slide_content(slide, i + 1)
                slides.append(slide_content)
            
            title = self._extract_presentation_title(prs)
            metadata = self._extract_metadata(prs)
            
            return PresentationData(
                title=title,
                slides=slides,
                total_slides=len(slides),
                metadata=metadata
            )
        
        except Exception as e:
            self.logger.error(f"Error reading PowerPoint file: {e}")
            raise
    
    def _extract_slide_content(self, slide: Slide, slide_number: int) -> SlideContent:
        """個別スライドの内容を抽出"""
        title = self._extract_slide_title(slide)
        text_content = self._extract_text_content(slide)
        slide_type = self._determine_slide_type(slide, slide_number)
        images = self._extract_images(slide)
        notes = self._extract_notes(slide)
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            slide_type=slide_type,
            images=images,
            notes=notes
        )
    
    def _extract_slide_title(self, slide: Slide) -> str:
        """スライドのタイトルを抽出"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # タイトルシェイプがない場合、最初のテキストボックスを探す
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return shape.text.strip().split('\n')[0]
        
        return f"スライド {slide.slide_id}"
    
    def _extract_text_content(self, slide: Slide) -> str:
        """スライドのテキスト内容をすべて抽出"""
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_text = self._extract_table_text(shape)
                if table_text:
                    text_parts.append(table_text)
        
        return '\n\n'.join(text_parts)
    
    def _extract_table_text(self, table_shape) -> str:
        """テーブルからテキストを抽出"""
        try:
            table = table_shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(' | '.join(cells))
            return '\n'.join(rows)
        except Exception as e:
            self.logger.warning(f"Error extracting table text: {e}")
            return ""
    
    def _determine_slide_type(self, slide: Slide, slide_number: int) -> SlideType:
        """スライドの種類を判定"""
        title = self._extract_slide_title(slide).lower()
        
        if slide_number == 1 or any(keyword in title for keyword in ['タイトル', 'title', '概要', 'overview']):
            return SlideType.TITLE
        elif any(keyword in title for keyword in ['まとめ', 'conclusion', '結論', '終わり']):
            return SlideType.CONCLUSION
        elif any(keyword in title for keyword in ['章', 'section', 'part', '第']):
            return SlideType.SECTION
        else:
            return SlideType.CONTENT
    
    def _extract_images(self, slide: Slide) -> List[str]:
        """スライド内の画像情報を抽出"""
        images = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_info = self._get_image_info(shape)
                    images.append(image_info)
                except Exception as e:
                    self.logger.warning(f"Error extracting image info: {e}")
                    images.append("画像（詳細取得エラー）")
        
        return images
    
    def _get_image_info(self, picture_shape) -> str:
        """画像の詳細情報を取得"""
        try:
            # 画像の基本情報
            width = picture_shape.width
            height = picture_shape.height
            
            # 画像の説明やaltテキストがあれば取得
            alt_text = getattr(picture_shape, 'alt_text', '')
            
            info = f"画像 ({width}x{height})"
            if alt_text:
                info += f" - {alt_text}"
            
            return info
        except Exception as e:
            self.logger.warning(f"Error getting image details: {e}")
            return "画像"
    
    def _extract_notes(self, slide: Slide) -> Optional[str]:
        """スライドノートを抽出"""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    return notes_slide.notes_text_frame.text.strip()
        except Exception as e:
            self.logger.warning(f"Error extracting notes: {e}")
        
        return None
    
    def _extract_presentation_title(self, prs: Presentation) -> str:
        """プレゼンテーション全体のタイトルを取得"""
        if prs.slides:
            first_slide_title = self._extract_slide_title(prs.slides[0])
            if first_slide_title and first_slide_title != f"スライド {prs.slides[0].slide_id}":
                return first_slide_title
        
        return "無題のプレゼンテーション"
    
    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """プレゼンテーションのメタデータを抽出"""
        try:
            core_props = prs.core_properties
            return {
                "author": core_props.author or "不明",
                "title": core_props.title or "",
                "subject": core_props.subject or "",
                "created": core_props.created.isoformat() if core_props.created else None,
                "modified": core_props.modified.isoformat() if core_props.modified else None,
                "last_modified_by": core_props.last_modified_by or "不明",
                "revision": core_props.revision,
                "slide_count": len(prs.slides)
            }
        except Exception as e:
            self.logger.warning(f"Error extracting metadata: {e}")
            return {"slide_count": len(prs.slides)}
    
    def get_slide_summary(self, presentation: PresentationData) -> str:
        """プレゼンテーション全体の要約を生成"""
        summary_parts = [
            f"タイトル: {presentation.title}",
            f"総スライド数: {presentation.total_slides}",
            f"作成者: {presentation.metadata.get('author', '不明')}"
        ]
        
        if presentation.slides:
            summary_parts.append("\nスライド構成:")
            for slide in presentation.slides:
                summary_parts.append(f"  {slide.slide_number}. {slide.title} ({slide.slide_type.value})")
        
        return '\n'.join(summary_parts)